# knowledge_flow_app/processors/pptx/pptx_processor.py

import logging
import shutil
from pathlib import Path

from pptx import Presentation
from knowledge_flow_app.input_processors.base_input_processor import BaseMarkdownProcessor

logger = logging.getLogger(__name__)

class PptxMarkdownProcessor(BaseMarkdownProcessor):
    def check_file_validity(self, file_path: Path) -> bool:
        """Checks if the PPTX file is valid and can be opened."""
        try:
            Presentation(str(file_path))
            return True
        except Exception as e:
            logger.error(f"Invalid or corrupted PPTX file: {file_path} - {e}")
            return False

    def extract_file_metadata(self, file_path: Path) -> dict:
        """Extracts basic metadata from the PPTX file."""
        metadata = {"document_name": file_path.name}
        try:
            presentation = Presentation(str(file_path))
            metadata["num_slides"] = len(presentation.slides)
        except Exception as e:
            logger.error(f"Error reading PPTX file: {e}")
            metadata["error"] = str(e)
        return metadata

    def convert_file_to_markdown(self, file_path: Path, output_dir: Path) -> dict:
        """Converts each slide's text content into Markdown."""

        md_path = output_dir / "output.md"

        try:
            presentation = Presentation(str(file_path))
            slide_texts = []

            for slide in presentation.slides:
                slide_md = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_md.append(shape.text.strip())
                if slide_md:
                    slide_texts.append("### Slide\n" + "\n\n".join(slide_md))

            content = "\n\n---\n\n".join(slide_texts) if slide_texts else "*No extractable text*"
            md_path.write_text(content)
            return {
                "doc_dir": str(output_dir),
                "md_file": str(md_path),
                "status": "success",
                "message": "PPTX slides converted to Markdown.",
            }

        except Exception as e:
            logger.error(f"Failed to convert PPTX to Markdown: {e}")
            return {
                "doc_dir": str(output_dir),
                "md_file": None,
                "status": "error",
                "message": str(e),
            }
